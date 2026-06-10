"""
Генератор презентации MAPS для руководства.
Запуск: .venv/bin/python docs/generate_presentation.py
Результат: docs/MAPS_Презентация.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Цвета ──────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2E)   # тёмно-синий фон
BG_SLIDE   = RGBColor(0x11, 0x22, 0x3A)   # чуть светлее для контента
BLUE_ACT   = RGBColor(0x29, 0x9D, 0xFF)   # акцентный голубой
BLUE_DARK  = RGBColor(0x16, 0x5C, 0xB0)   # тёмно-синий акцент
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)   # оранжевый (проблемы)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)   # зелёный (результат)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xB0, 0xC4, 0xDE)
GRAY_MID   = RGBColor(0x5A, 0x7A, 0x9A)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)

# ── Размер слайда (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    """Добавляет полностью пустой слайд."""
    layout = prs.slide_layouts[6]   # Blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = BG_DARK):
    """Заливает фон слайда цветом."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def txb(slide, text, x, y, w, h,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Добавляет текстовый блок."""
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def txb_multi(slide, lines, x, y, w, h,
              size=20, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Добавляет текстовый блок с несколькими параграфами."""
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line == "":
            p.add_run().text = ""
            continue
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def accent_bar(slide, color=BLUE_ACT, y=Inches(0.55), height=Inches(0.06)):
    """Горизонтальная цветная полоса под заголовком."""
    add_rect(slide, Inches(0.5), y, Inches(12.33), height, color)


# ══════════════════════════════════════════════════════════════════════════
#  СЛАЙДЫ
# ══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Титульный слайд."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    # Крупный фоновый прямоугольник-акцент слева
    add_rect(s, Inches(0), Inches(0), Inches(0.25), H, BLUE_ACT)

    # Полоса по центру
    add_rect(s, Inches(0.5), Inches(3.5), Inches(6.5), Inches(0.05), BLUE_ACT)

    txb(s, "MAPS", Inches(0.6), Inches(1.0), Inches(10), Inches(1.8),
        size=96, bold=True, color=BLUE_ACT, align=PP_ALIGN.LEFT)

    txb(s, "Material Allocation & Planning System",
        Inches(0.6), Inches(2.7), Inches(10), Inches(0.7),
        size=22, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

    txb(s, "Система управления материальным обеспечением работ",
        Inches(0.6), Inches(3.7), Inches(10), Inches(0.8),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txb(s, "Как перестать планировать вслепую\nи получить точную картину за секунды",
        Inches(0.6), Inches(4.5), Inches(9), Inches(1.2),
        size=20, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

    # Год и компания
    txb(s, "2026", Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
        size=16, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_02_question(prs):
    """Вопрос-крючок."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    # Большая цитата по центру
    txb(s,
        "«Какие работы из нашего плана\nобеспечены материалами —\nа какие нет,\nи на сколько именно не хватает?»",
        Inches(1.0), Inches(0.8), Inches(11.33), Inches(4.5),
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(4.0), Inches(5.5), Inches(5.33), Inches(0.05), BLUE_ACT)

    txb(s, "Прямо сейчас — кто может ответить на этот вопрос?",
        Inches(1.0), Inches(5.7), Inches(11.33), Inches(0.8),
        size=22, bold=False, color=GRAY_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    txb(s, "02", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_03_scale(prs):
    """Масштаб задачи."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Масштаб, при котором ручное управление не работает",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s)

    # Три блока с цифрами
    boxes = [
        ("30+ млрд ₸",   "ежегодный объём\nстроительно-монтажных\nи ремонтных работ"),
        ("10 000+",       "позиций работ\nв годовой заявке\nпо всей компании"),
        ("10+ филиалов",  "каждый — несколько ЛПУ\nи множество складов\nс разными остатками"),
    ]
    bw = Inches(3.7)
    gap = Inches(0.5)
    for i, (num, desc) in enumerate(boxes):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(1.2), bw, Inches(3.8), BG_SLIDE)
        add_rect(s, bx, Inches(1.2), bw, Inches(0.08), BLUE_ACT)
        txb(s, num, bx, Inches(1.5), bw, Inches(1.4),
            size=48, bold=True, color=BLUE_ACT, align=PP_ALIGN.CENTER)
        txb(s, desc, bx, Inches(2.9), bw, Inches(1.8),
            size=18, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.5), Inches(5.3), Inches(12.33), Inches(0.05), ORANGE)
    txb(s, "При таком масштабе каждый процент неэффективности = сотни миллионов тенге",
        Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.7),
        size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    txb(s, "03", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_04_problems_a(prs):
    """Проблемы 1–3."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Что происходит сегодня",
        Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=ORANGE)

    problems = [
        ("01", "Нет единой картины обеспечённости",
         "Планировщики вручную сверяют таблицы остатков\nс потребностями по каждому филиалу.\nЭто занимает дни. Данные устаревают быстрее,\nчем успевают быть обработаны."),
        ("02", "Двойное распределение материала",
         "Один и тот же материал «запланирован» сразу\nна несколько работ. Одна получает материал,\nдругая — нет. Обе в плане считались обеспеченными.\nИтог: срыв срока, аварийная закупка, претензии."),
        ("03", "Аварийные работы ломают весь план",
         "Появляется новая аварийная работа —\nи плановая картина мгновенно устаревает.\nПересчитать вручную — снова дни.\nВсё это время работают с неактуальными данными."),
    ]

    bw = Inches(3.8)
    gap = Inches(0.35)
    for i, (num, title, desc) in enumerate(problems):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(1.1), bw, Inches(5.5), BG_SLIDE)
        add_rect(s, bx, Inches(1.1), bw, Inches(0.08), ORANGE)
        txb(s, num, bx + Inches(0.2), Inches(1.25), Inches(0.8), Inches(0.7),
            size=36, bold=True, color=ORANGE)
        txb(s, title, bx + Inches(0.2), Inches(1.95), bw - Inches(0.4), Inches(0.9),
            size=17, bold=True, color=WHITE)
        txb(s, desc, bx + Inches(0.2), Inches(2.9), bw - Inches(0.4), Inches(3.2),
            size=15, color=GRAY_LIGHT)

    txb(s, "04", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_05_problems_b(prs):
    """Проблемы 4–5 (деньги)."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Где теряются деньги",
        Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=ORANGE)

    # Проблема 4
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(2.65), BG_SLIDE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(0.08), Inches(2.65), ORANGE)
    txb(s, "04  Закупают не то и не столько",
        Inches(0.75), Inches(1.2), Inches(11.5), Inches(0.6),
        size=20, bold=True, color=ORANGE)
    txb_multi(s, [
        "· Закупают то, что уже есть на складе соседнего филиала",
        "· Не учитывают законтрактованные поставки «в пути»",
        "· Не учитывают фактические списания по уже выполненным объёмам",
        "Результат: избыток в одних местах — острый дефицит в других",
    ], Inches(0.85), Inches(1.85), Inches(11.5), Inches(1.7), size=17, color=GRAY_LIGHT)

    # Проблема 5
    add_rect(s, Inches(0.5), Inches(4.0), Inches(12.33), Inches(2.0), BG_SLIDE)
    add_rect(s, Inches(0.5), Inches(4.0), Inches(0.08), Inches(2.0), ORANGE)
    txb(s, "05  Стоимость работ определяется постфактум",
        Inches(0.75), Inches(4.1), Inches(11.5), Inches(0.6),
        size=20, bold=True, color=ORANGE)
    txb(s, "Фактическая стоимость каждой работы становится известна после того,\n"
           "как повлиять на ситуацию уже невозможно. Экономия или перерасход — сюрприз.",
        Inches(0.85), Inches(4.75), Inches(11.5), Inches(1.0),
        size=17, color=GRAY_LIGHT)

    txb(s, "05", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_06_bridge(prs):
    """Переход к решению."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    add_rect(s, Inches(0), Inches(0), Inches(0.25), H, BLUE_ACT)

    txb(s, "Всё, что описано выше —\nне константа.",
        Inches(1.0), Inches(1.5), Inches(11), Inches(2.0),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txb(s, "Это задача автоматизации.",
        Inches(1.0), Inches(3.5), Inches(11), Inches(1.0),
        size=36, bold=True, color=BLUE_ACT, align=PP_ALIGN.CENTER)

    txb(s, "Она решена.",
        Inches(1.0), Inches(4.5), Inches(11), Inches(1.0),
        size=36, bold=False, color=GREEN, align=PP_ALIGN.CENTER, italic=True)

    txb(s, "06", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_07_what_is_maps(prs):
    """Что такое MAPS."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "MAPS — автоматическое распределение материалов в реальном времени",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.8),
        size=24, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    # Главный тезис
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(1.1), BLUE_DARK)
    txb(s, "Система принимает данные из SAP и Excel — и за секунды выдаёт полную картину\n"
           "обеспечённости всех работ по всем филиалам.",
        Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.9),
        size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Три колонки — для кого
    roles = [
        ("Руководитель", "Картина обеспечённости\nв деньгах и %\nпо каждой работе"),
        ("Снабженец",    "Точная заявка\nна закупку\nбез двойного счёта"),
        ("Производственник", "Ответ на вопрос\n«мои работы обеспечены?»\nодним кликом"),
    ]
    bw = Inches(3.7)
    gap = Inches(0.5)
    for i, (role, desc) in enumerate(roles):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(2.5), bw, Inches(3.8), BG_SLIDE)
        add_rect(s, bx, Inches(2.5), bw, Inches(0.07), GREEN)
        txb(s, role, bx, Inches(2.65), bw, Inches(0.65),
            size=21, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txb(s, desc, bx, Inches(3.35), bw, Inches(2.5),
            size=18, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

    txb(s, "07", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_08_algorithm(prs):
    """Пятислойный алгоритм."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Пятислойный алгоритм — учитывает все источники покрытия",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=26, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    layers = [
        (1, "Фактические списания SAP",        "Уже проведено документально — учитывается первым",          GREEN),
        (2, "Выдано, но не списано",            "Физически на объекте, документ ещё не оформлен",            GREEN),
        (3, "Складские остатки",                "Свой завод → свой филиал → другие филиалы",                 BLUE_ACT),
        (4, "Законтрактованные поставки",       "Подтверждённые договоры — материалы «в пути»",              BLUE_ACT),
        (5, "ДЕФИЦИТ — точная цифра к закупу", "Остаток, который реально нужно купить — без двойного счёта", ORANGE),
    ]

    lh = Inches(0.98)
    ly_start = Inches(1.1)
    for i, (num, title, desc, color) in enumerate(layers):
        ly = ly_start + i * (lh + Inches(0.05))
        add_rect(s, Inches(0.5), ly, Inches(12.33), lh, BG_SLIDE)
        add_rect(s, Inches(0.5), ly, Inches(0.55), lh, color)
        txb(s, str(num), Inches(0.5), ly, Inches(0.55), lh,
            size=26, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        txb(s, title, Inches(1.2), ly + Inches(0.05), Inches(4.5), lh - Inches(0.1),
            size=18, bold=True, color=WHITE)
        txb(s, desc,  Inches(6.0), ly + Inches(0.05), Inches(6.6), lh - Inches(0.1),
            size=16, color=GRAY_LIGHT)

    txb(s, "Материал назначается ровно один раз. Дублирования нет.",
        Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5),
        size=17, bold=True, color=BLUE_ACT, align=PP_ALIGN.CENTER, italic=True)

    txb(s, "08", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_09_priority(prs):
    """Приоритизация работ."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Умная очередь: кто получает материал первым",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    priorities = [
        ("1", "Аварийные работы",  "Всегда получают материалы первыми — независимо от дат и других условий", ORANGE),
        ("2", "По дате начала",    "Чем раньше старт работы — тем выше место в очереди распределения",       BLUE_ACT),
        ("3", "По номеру приоритета", "Настраиваемый порядок внутри одной даты под задачи компании",         BLUE_DARK),
    ]

    for i, (num, title, desc, color) in enumerate(priorities):
        by = Inches(1.2) + i * Inches(1.55)
        add_rect(s, Inches(0.5), by, Inches(12.33), Inches(1.35), BG_SLIDE)
        add_rect(s, Inches(0.5), by, Inches(0.7), Inches(1.35), color)
        txb(s, num, Inches(0.5), by, Inches(0.7), Inches(1.35),
            size=32, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        txb(s, title, Inches(1.35), by + Inches(0.12), Inches(4.0), Inches(0.6),
            size=20, bold=True, color=WHITE)
        txb(s, desc, Inches(1.35), by + Inches(0.65), Inches(10.8), Inches(0.65),
            size=16, color=GRAY_LIGHT)

    add_rect(s, Inches(0.5), Inches(5.95), Inches(12.33), Inches(0.08), GREEN)
    txb(s, "Добавили аварийную работу → нажали «Пересчитать» → через секунды актуальная картина.\n"
           "Не через день. Не через два. Через секунды.",
        Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.9),
        size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    txb(s, "09", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_10_input_data(prs):
    """Входные данные."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Входные данные — то, что у вас уже есть",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    rows = [
        ("Перечень работ",             "Коды, наименования, сроки, приоритеты, филиал, завод",      "Производственный план / SAP"),
        ("Аварийные работы",           "Работы наивысшего приоритета, добавляемые в ходе периода",  "Производство / диспетчерская"),
        ("Потребности в материалах",   "Какой материал нужен для каждой работы и в каком кол-ве",   "Отдел снабжения / SAP"),
        ("Складские остатки",          "Партии: количество, цена, дата поступления, склад",          "SAP / складской учёт"),
        ("Законтракт. поставки",       "Материалы в пути: договор, поставщик, кол-во, дата",         "Отдел снабжения / SAP"),
        ("Фактические списания",       "Что уже проведено документально на работы в SAP",            "SAP (уже интегрировано ✓)"),
        ("Выдано не списано",          "Выдано со склада — документ ещё не проведён",                "Кладовщики / SAP"),
    ]

    rh = Inches(0.72)
    col_w = [Inches(2.6), Inches(5.7), Inches(3.7)]
    col_x = [Inches(0.5), Inches(3.2), Inches(9.0)]

    # Заголовок таблицы
    for ci, (cw, cx, ctitle) in enumerate(zip(col_w, col_x,
            ["Источник данных", "Что содержит", "Откуда берётся"])):
        add_rect(s, cx, Inches(1.1), cw - Inches(0.05), Inches(0.5), BLUE_DARK)
        txb(s, ctitle, cx + Inches(0.1), Inches(1.15), cw - Inches(0.2), Inches(0.4),
            size=14, bold=True, color=WHITE)

    for ri, row in enumerate(rows):
        ry = Inches(1.65) + ri * rh
        row_bg = BG_SLIDE if ri % 2 == 0 else RGBColor(0x14, 0x28, 0x45)
        for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
            add_rect(s, cx, ry, cw - Inches(0.05), rh - Inches(0.04), row_bg)
            clr = GREEN if "интегрировано" in val else (WHITE if ci == 0 else GRAY_LIGHT)
            txb(s, val, cx + Inches(0.1), ry + Inches(0.1),
                cw - Inches(0.2), rh - Inches(0.15),
                size=13, color=clr, bold=(ci == 0))

    txb(s, "10", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_11_output(prs):
    """Что система выдаёт."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Результаты — три формата для разных задач",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=GREEN)

    # Блок 1: Excel
    add_rect(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(5.5), BG_SLIDE)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(0.07), GREEN)
    txb(s, "Excel-отчёт  (7 листов)",
        Inches(0.65), Inches(1.2), Inches(5.6), Inches(0.55),
        size=18, bold=True, color=GREEN)
    sheets = [
        ("Распределение",        "каждая работа × материал × источник × стоимость"),
        ("Обеспеченность",       "% и стоимость покрытия по каждой работе"),
        ("Движение склада",      "что взято, откуда, по какой цене"),
        ("Остатки складов",      "что осталось после распределения"),
        ("Возм. перемещение",    "где взять у соседнего филиала вместо закупки"),
        ("Поставки в работу",    "что «в пути» уже распределено на работы"),
        ("Заявка на закупку",    "готовый список дефицита — к отделу снабжения"),
    ]
    for i, (name, desc) in enumerate(sheets):
        sy = Inches(1.85) + i * Inches(0.66)
        txb(s, f"  {name}", Inches(0.65), sy, Inches(2.1), Inches(0.55),
            size=13, bold=True, color=WHITE)
        txb(s, desc, Inches(2.8), sy, Inches(3.5), Inches(0.55),
            size=12, color=GRAY_LIGHT)

    # Блок 2: Дашборд
    add_rect(s, Inches(6.7), Inches(1.1), Inches(6.1), Inches(2.5), BG_SLIDE)
    add_rect(s, Inches(6.7), Inches(1.1), Inches(6.1), Inches(0.07), BLUE_ACT)
    txb(s, "Веб-дашборд в реальном времени",
        Inches(6.85), Inches(1.2), Inches(5.8), Inches(0.55),
        size=18, bold=True, color=BLUE_ACT)
    txb_multi(s, [
        "· Сводная статистика по всем работам и складам",
        "· Запуск перерасчёта одной кнопкой",
        "· Доступен с любого устройства через браузер",
        "· Без установки ПО на рабочие места",
    ], Inches(6.85), Inches(1.85), Inches(5.8), Inches(1.5), size=14, color=GRAY_LIGHT)

    # Блок 3: Заявка
    add_rect(s, Inches(6.7), Inches(3.8), Inches(6.1), Inches(2.8), BG_SLIDE)
    add_rect(s, Inches(6.7), Inches(3.8), Inches(6.1), Inches(0.07), ORANGE)
    txb(s, "Точная заявка на закупку",
        Inches(6.85), Inches(3.9), Inches(5.8), Inches(0.55),
        size=18, bold=True, color=ORANGE)
    txb_multi(s, [
        "Автоматически рассчитанный список дефицитных",
        "позиций с количеством и стоимостью.",
        "",
        "Исключает двойной счёт.",
        "Не требует ручной сверки.",
    ], Inches(6.85), Inches(4.55), Inches(5.8), Inches(1.8), size=15, color=GRAY_LIGHT)

    txb(s, "11", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_12_ai(prs):
    """AI-возможности."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Встроенный AI: не просто цифры — объяснения и прогнозы",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=26, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    features = [
        ("Анализ дефицита", "Уже работает",
         "Для каждой дефицитной позиции AI объясняет на русском языке:\n"
         "почему не хватило · что покрыл каждый слой\n"
         "откуда брать · сколько и когда купить\n\n"
         "Готовое объяснение для руководителя — без разбора таблиц.",
         GREEN),
        ("Прогнозирование дефицитов", "В разработке",
         "На основе динамики списаний и темпов выполнения\n"
         "система предупреждает о дефиците до того,\n"
         "как он возникнет.",
         BLUE_ACT),
        ("Рекомендации по оптимизации", "В разработке",
         "Что выгоднее: переместить между филиалами\n"
         "или закупить. Где риск срыва срока.\n"
         "Какие работы можно начать раньше.",
         BLUE_ACT),
    ]

    bw = Inches(3.8)
    gap = Inches(0.35)
    for i, (title, status, desc, color) in enumerate(features):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(1.1), bw, Inches(5.5), BG_SLIDE)
        add_rect(s, bx, Inches(1.1), bw, Inches(0.07), color)
        txb(s, title, bx + Inches(0.2), Inches(1.25), bw - Inches(0.4), Inches(0.65),
            size=18, bold=True, color=WHITE)
        # Бейдж статуса
        add_rect(s, bx + Inches(0.2), Inches(1.95), Inches(1.6), Inches(0.38),
                 color if status == "Уже работает" else GRAY_MID)
        txb(s, status, bx + Inches(0.2), Inches(1.97), Inches(1.6), Inches(0.35),
            size=11, bold=True, color=BG_DARK if status == "Уже работает" else WHITE,
            align=PP_ALIGN.CENTER)
        txb(s, desc, bx + Inches(0.2), Inches(2.5), bw - Inches(0.4), Inches(3.8),
            size=15, color=GRAY_LIGHT)

    txb(s, "12", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_13_before_after(prs):
    """До и после."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Что меняется с MAPS",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.6),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=GREEN)

    rows = [
        ("Картина обеспечённости",       "Дни ручной работы",                    "Секунды"),
        ("Точность заявки на закупку",   "Двойной счёт, пропущенные поставки",   "Точный расчёт по всем источникам"),
        ("Реакция на аварийную работу",  "Перепланирование вручную, 1–2 дня",    "Мгновенный пересчёт"),
        ("Межфилиальное перемещение",    "Не анализируется систематически",       "Автоматическая подсказка"),
        ("Стоимость работ",              "Постфактум — уже нельзя повлиять",     "При каждом распределении"),
        ("Прозрачность решений",         "Сложно восстановить кто что решил",    "Полный аудиторский след"),
    ]

    # Заголовки колонок
    col_conf = [(Inches(0.5), Inches(4.3)), (Inches(5.0), Inches(4.0)), (Inches(9.2), Inches(4.0))]
    col_titles = ["", "До MAPS", "С MAPS"]
    col_colors = [GRAY_MID, ORANGE, GREEN]
    for (cx, cw), ct, cc in zip(col_conf, col_titles, col_colors):
        if ct:
            add_rect(s, cx, Inches(1.05), cw - Inches(0.1), Inches(0.48), cc)
            txb(s, ct, cx, Inches(1.1), cw - Inches(0.1), Inches(0.4),
                size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    rh = Inches(0.82)
    for ri, (aspect, before, after) in enumerate(rows):
        ry = Inches(1.6) + ri * rh
        bg_row = BG_SLIDE if ri % 2 == 0 else RGBColor(0x14, 0x28, 0x45)
        for (cx, cw), val, is_after in zip(col_conf,
                [aspect, before, after],
                [False, False, True]):
            add_rect(s, cx, ry, cw - Inches(0.1), rh - Inches(0.06), bg_row)
            clr = GREEN if is_after else (WHITE if not is_after and val == aspect else ORANGE)
            txb(s, val, cx + Inches(0.15), ry + Inches(0.1),
                cw - Inches(0.3), rh - Inches(0.2),
                size=14, color=clr, bold=(val == aspect))

    txb(s, "13", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_14_tech(prs):
    """Технологический стек."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Технологии корпоративного класса — не самописные «костыли»",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=24, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    # Блок цифр
    metrics = [
        ("~8 000",  "строк кода"),
        ("15",      "сущностей\nв БД"),
        ("9",       "модулей\nимпорта"),
        ("21",      "шаг движка\nраспределения"),
    ]
    mw = Inches(2.8)
    for i, (num, label) in enumerate(metrics):
        mx = Inches(0.5) + i * (mw + Inches(0.24))
        add_rect(s, mx, Inches(1.1), mw, Inches(1.55), BG_SLIDE)
        add_rect(s, mx, Inches(1.1), mw, Inches(0.07), BLUE_ACT)
        txb(s, num, mx, Inches(1.2), mw, Inches(0.8),
            size=36, bold=True, color=BLUE_ACT, align=PP_ALIGN.CENTER)
        txb(s, label, mx, Inches(1.95), mw, Inches(0.6),
            size=13, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # Таблица стека
    tech_rows = [
        ("Python",             "Язык №1 для аналитических систем",               "Google · NASA · крупнейшие банки"),
        ("FastAPI",            "Один из самых быстрых серверных фреймворков",      "Microsoft · Uber · Netflix"),
        ("PostgreSQL",         "Промышленная СУБД мирового класса",               "Финансовые реестры, медицинские системы"),
        ("pandas",             "Индустриальный стандарт обработки данных",        "Банки · биржи · крупные производства"),
        ("Pydantic v2",        "Автоматическая проверка каждой строки данных",    "Ошибки отсеиваются до записи в базу"),
        ("SQLAlchemy + Alembic","ORM и версионирование схемы БД",                 "Изменения применяются безопасно, как в банке"),
        ("Claude AI (Anthropic)","Самая продвинутая коммерческая языковая модель","Та же технология — лучшие корпоративные AI"),
    ]

    col_w = [Inches(2.5), Inches(4.5), Inches(4.8)]
    col_x = [Inches(0.5), Inches(3.15), Inches(7.75)]
    rh = Inches(0.57)

    for ci, (cw, cx, ct) in enumerate(zip(col_w, col_x, ["Технология", "Что это даёт", "Кто ещё использует"])):
        add_rect(s, cx, Inches(2.85), cw - Inches(0.07), Inches(0.42), BLUE_DARK)
        txb(s, ct, cx + Inches(0.1), Inches(2.88), cw - Inches(0.2), Inches(0.38),
            size=13, bold=True, color=WHITE)

    for ri, row in enumerate(tech_rows):
        ry = Inches(3.33) + ri * rh
        row_bg = BG_SLIDE if ri % 2 == 0 else RGBColor(0x14, 0x28, 0x45)
        for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
            add_rect(s, cx, ry, cw - Inches(0.07), rh - Inches(0.04), row_bg)
            clr = BLUE_ACT if ci == 0 else (WHITE if ci == 1 else GRAY_LIGHT)
            txb(s, val, cx + Inches(0.1), ry + Inches(0.07),
                cw - Inches(0.2), rh - Inches(0.1),
                size=13, bold=(ci == 0), color=clr)

    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.06), BLUE_ACT)
    txb(s, "Каждый инструмент выбран как лучшее решение для задач именно такого масштаба.",
        Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.5),
        size=15, bold=True, color=BLUE_ACT, align=PP_ALIGN.CENTER, italic=True)

    txb(s, "14", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_15_roadmap(prs):
    """Дорожная карта."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Дорожная карта развития",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=28, bold=True, color=WHITE)
    accent_bar(s, color=BLUE_ACT)

    phases = [
        ("Сейчас", "Готово к работе",
         ["Импорт данных из Excel и SAP",
          "Пятислойный алгоритм распределения",
          "AI-анализ дефицитов",
          "Веб-дашборд и Excel-отчёты",
          "Полный аудиторский след"],
         GREEN),
        ("6–12 месяцев", "Следующий этап",
         ["АРМ снабженца",
          "АРМ производственника",
          "Уведомления и согласование заявок",
          "Мобильный доступ",
          "История заявок и статусов поставок"],
         BLUE_ACT),
        ("12–24 месяца", "Стратегическое развитие",
         ["Прямой обмен данными с SAP",
          "Прогнозирование дефицитов",
          "Оптимизация межфилиальных перемещений",
          "Сводная аналитика по компании",
          "Анализ критического пути работ"],
         BLUE_DARK),
    ]

    bw = Inches(3.8)
    gap = Inches(0.45)
    for i, (period, status, items, color) in enumerate(phases):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(1.1), bw, Inches(5.5), BG_SLIDE)
        add_rect(s, bx, Inches(1.1), bw, Inches(0.08), color)
        txb(s, period, bx + Inches(0.2), Inches(1.25), bw - Inches(0.4), Inches(0.6),
            size=20, bold=True, color=color)
        txb(s, status, bx + Inches(0.2), Inches(1.85), bw - Inches(0.4), Inches(0.45),
            size=14, color=GRAY_LIGHT, italic=True)
        for j, item in enumerate(items):
            txb(s, f"✓  {item}" if i == 0 else f"·  {item}",
                bx + Inches(0.2), Inches(2.45) + j * Inches(0.72),
                bw - Inches(0.4), Inches(0.65),
                size=14, color=WHITE if i == 0 else GRAY_LIGHT)

    txb(s, "Каждый следующий шаг строится на уже работающем фундаменте. Переписывать ничего не нужно.",
        Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.45),
        size=15, color=GRAY_MID, align=PP_ALIGN.CENTER, italic=True)

    txb(s, "15", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_16_finale(prs):
    """Финал — возврат к вопросу."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    add_rect(s, Inches(0), Inches(0), Inches(0.25), H, GREEN)

    txb(s, "«Какие работы обеспечены материалами —\nа какие нет, и на сколько?»",
        Inches(0.6), Inches(0.6), Inches(12.0), Inches(1.8),
        size=30, bold=True, color=GRAY_LIGHT, align=PP_ALIGN.LEFT)

    add_rect(s, Inches(0.6), Inches(2.55), Inches(11.5), Inches(0.06), GREEN)

    txb(s, "С MAPS — этот вопрос получает ответ за секунды.",
        Inches(0.6), Inches(2.75), Inches(12.0), Inches(0.75),
        size=28, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

    details = [
        "По каждой из 10 000+ работ",
        "По всем 10+ филиалам",
        "С разбивкой по источникам, стоимости и % обеспечённости",
        "В любой момент — при любом изменении данных",
    ]
    for i, line in enumerate(details):
        txb(s, f"  →  {line}",
            Inches(0.6), Inches(3.65) + i * Inches(0.62),
            Inches(12.0), Inches(0.55),
            size=20, color=WHITE)

    txb(s, "16", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def slide_17_cta(prs):
    """Призыв к действию."""
    s = blank_slide(prs)
    bg(s, BG_DARK)

    txb(s, "Следующий шаг",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.7),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    accent_bar(s, color=BLUE_ACT, y=Inches(0.85))

    options = [
        ("01", "Пилотный запуск",
         "Один филиал, реальные данные следующего квартала.\nРезультат — за одну рабочую сессию.",
         GREEN),
        ("02", "Расширенная демонстрация",
         "На данных любого выбранного периода.\nПоказываем весь цикл: загрузка → расчёт → отчёт.",
         BLUE_ACT),
        ("03", "Рабочая встреча",
         "Детальное обсуждение с руководителями\nснабжения и производства.",
         ORANGE),
    ]

    bw = Inches(3.8)
    gap = Inches(0.45)
    for i, (num, title, desc, color) in enumerate(options):
        bx = Inches(0.5) + i * (bw + gap)
        add_rect(s, bx, Inches(1.3), bw, Inches(4.5), BG_SLIDE)
        add_rect(s, bx, Inches(1.3), bw, Inches(0.08), color)
        txb(s, num, bx + Inches(0.2), Inches(1.45), Inches(0.8), Inches(0.8),
            size=40, bold=True, color=color)
        txb(s, title, bx + Inches(0.2), Inches(2.3), bw - Inches(0.4), Inches(0.7),
            size=20, bold=True, color=WHITE)
        txb(s, desc, bx + Inches(0.2), Inches(3.1), bw - Inches(0.4), Inches(2.4),
            size=16, color=GRAY_LIGHT)

    add_rect(s, Inches(0.5), Inches(6.05), Inches(12.33), Inches(0.08), BLUE_ACT)
    txb(s, "MAPS · Material Allocation & Planning System · 2026",
        Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.55),
        size=16, color=GRAY_MID, align=PP_ALIGN.CENTER)

    txb(s, "17", Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
        size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
#  СБОРКА
# ══════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_question(prs)
    slide_03_scale(prs)
    slide_04_problems_a(prs)
    slide_05_problems_b(prs)
    slide_06_bridge(prs)
    slide_07_what_is_maps(prs)
    slide_08_algorithm(prs)
    slide_09_priority(prs)
    slide_10_input_data(prs)
    slide_11_output(prs)
    slide_12_ai(prs)
    slide_13_before_after(prs)
    slide_14_tech(prs)
    slide_15_roadmap(prs)
    slide_16_finale(prs)
    slide_17_cta(prs)

    out = "docs/MAPS_Презентация.pptx"
    prs.save(out)
    print(f"✓ Сохранено: {out}  ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()
